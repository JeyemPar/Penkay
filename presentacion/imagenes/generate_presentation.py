import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Widescreen 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    base_dir = r"D:\Terraventure_Hackaton\penkay\presentacion\imagenes"
    
    # Colors based on Penkay identity
    # verde bosque (forest green): #228B22 -> 34, 139, 34
    # oliva (olive): #808000 -> 128, 128, 0
    # lima (lime): #32CD32 -> 50, 205, 50
    # ocre (ochre): #CC7722 -> 204, 119, 34
    # crema (cream): #FFFDD0 -> 255, 253, 208

    color_forest = RGBColor(34, 139, 34)
    color_dark_forest = RGBColor(20, 80, 20)
    color_cream = RGBColor(255, 253, 208)
    color_ochre = RGBColor(204, 119, 34)

    def add_overlay_slide(image_filename, title, subtitle="", text_pos="left"):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        img_path = os.path.join(base_dir, image_filename)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        if text_pos == "left":
            left = Inches(0.5)
            top = Inches(2.5)
            width = Inches(5.5)
            height = Inches(3)
        elif text_pos == "bottom":
            left = Inches(0)
            top = Inches(6.0)
            width = prs.slide_width
            height = Inches(1.5)
        elif text_pos == "top":
            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = Inches(1.2)
        else:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(12)
            height = Inches(1.5)

        # Add a semi-transparent or solid background shape for text readability
        if text_pos in ["bottom", "top"]:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color_dark_forest
            # shape.fill.transparency = 0.2 # pptx handles transparency poorly sometimes, solid is safer
            shape.line.fill.background()
            
            txBox = slide.shapes.add_textbox(left + Inches(0.5), top, width - Inches(1), height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(36)
            p.font.color.rgb = color_cream
            p.alignment = PP_ALIGN.CENTER
            
            if subtitle:
                p2 = tf.add_paragraph()
                p2.text = subtitle
                p2.font.size = Pt(24)
                p2.font.color.rgb = color_cream
                p2.alignment = PP_ALIGN.CENTER

        elif text_pos == "left":
            # Just text box with shadow or bold colors
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(54)
            p.font.color.rgb = color_cream
            
            if subtitle:
                p2 = tf.add_paragraph()
                p2.text = subtitle
                p2.font.size = Pt(32)
                p2.font.color.rgb = color_cream

    # Slide 1: Portada
    add_overlay_slide(
        "01-portada-penkay-campo-satelite.png",
        "PENKAY",
        "Inteligencia Ambiental y\nTecnológica para el Campo",
        text_pos="left"
    )

    # Slide 2: Arquitectura
    add_overlay_slide(
        "02-arquitectura-tecnologica-penkay.png",
        "Arquitectura Tecnológica",
        "Integración de Sentinel, captura móvil, IA y trazabilidad web",
        text_pos="bottom"
    )

    # Slide 3: Machine Learning
    add_overlay_slide(
        "03-machine-learning-vision-artificial.png",
        "Visión Artificial y Machine Learning",
        "Captura ➔ Preparación ➔ Visión Artificial ➔ Decisión",
        text_pos="bottom"
    )

    # Slide 4: Impacto Ambiental
    add_overlay_slide(
        "04-impacto-ambiental-medible.png",
        "Impacto Ambiental Medible",
        "NDVI, Humedad, Carbono Orgánico del Suelo, Erosión y Huella Hash",
        text_pos="bottom"
    )

    # Slide 5: Drones
    add_overlay_slide(
        "05-drones-calibracion-multiespectral.png",
        "Drones para Calibración",
        "Levantamientos periódicos de alta resolución para calibrar observaciones satelitales",
        text_pos="bottom"
    )

    # Slide 6: Academia, Comunidad, Productores
    add_overlay_slide(
        "06-academia-comunidad-productores-campo.png",
        "Trabajo de Campo Colaborativo",
        "Sinergia horizontal: comunidad, academia y jóvenes PencoTech",
        text_pos="bottom"
    )

    # Slide 7: Ecosistema
    add_overlay_slide(
        "07-ecosistema-comunidad-academia-productores.png",
        "Ecosistema Penkay",
        "Relación bidireccional: Datos de campo ➔ Conocimiento ➔ Decisiones",
        text_pos="bottom"
    )

    output_path = os.path.join(base_dir, "Presentacion_Penkay.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
